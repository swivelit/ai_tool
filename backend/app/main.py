from __future__ import annotations

import json
import os
import re
import sys
import tempfile
import time
from datetime import date, datetime, timedelta
from pathlib import Path
from typing import Any, Dict, List, Optional
from zoneinfo import ZoneInfo

from dotenv import load_dotenv
from fastapi import Depends, FastAPI, File, HTTPException, Query, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from openai import OpenAI
from pydantic import BaseModel
from sqlalchemy import inspect as sa_inspect, text as sa_text
from sqlmodel import Session, delete, select

from .database import engine, get_session
from .models import Conversation, DailyRoutine, Item, QACache, User, UserProfile

CURRENT_DIR = Path(__file__).resolve().parent
BACKEND_ROOT = CURRENT_DIR.parent
if str(BACKEND_ROOT) not in sys.path:
    sys.path.insert(0, str(BACKEND_ROOT))

from config import LOGS_DIR, PIPELINE_VERSION  # noqa: E402
from stage_behaviour_questions import BehaviourQuestionnaire, QUESTIONS as PIPELINE_QUESTIONS  # noqa: E402
from stage_english_remodel import EnglishRemodeler  # noqa: E402
from stage_openai_core import OpenAICore  # noqa: E402
from stage_translate import StageTranslator  # noqa: E402

load_dotenv()

PERSONALITY_QUESTIONS_VERSION = 1
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
OPENAI_JSON_MODEL = os.getenv("OPENAI_JSON_MODEL", "gpt-4o-mini")
if not OPENAI_API_KEY:
    raise RuntimeError("OPENAI_API_KEY is not set (env var missing)")

client = OpenAI(api_key=OPENAI_API_KEY)

app = FastAPI(title="J AI Backend")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

STAGE_BEHAVIOUR = BehaviourQuestionnaire()
STAGE_CORE = OpenAICore()
STAGE_REMODELER = EnglishRemodeler(STAGE_CORE)
STAGE_TRANSLATOR = StageTranslator(STAGE_CORE)
STAGE_CACHE: Dict[str, Dict[str, Any]] = {}


def _normalize_reply_language(value: Optional[str]) -> str:
    normalized = str(value or "").strip().lower()
    if normalized in {"en", "english"}:
        return "en"
    if normalized in {"ta", "tamil", "mixed", "tanglish"}:
        return "ta"
    return "ta"


def _normalize_lookup_text(text: str) -> str:
    parts = re.findall(r"[a-z0-9_\u0B80-\u0BFF]+", str(text or "").lower())
    return " ".join(parts)


def _token_overlap_score(left: str, right: str) -> float:
    left_tokens = set(_normalize_lookup_text(left).split())
    right_tokens = set(_normalize_lookup_text(right).split())
    if not left_tokens and not right_tokens:
        return 1.0
    if not left_tokens or not right_tokens:
        return 0.0
    return len(left_tokens & right_tokens) / max(1, len(left_tokens | right_tokens))


def _build_pipeline_result(
    *,
    raw_english: str,
    remodeled_english: Optional[str] = None,
    tamil_text: str = "",
    theni_tamil_text: str = "",
    route_taken: str,
    direct_answer_source: str = "",
    direct_answer_confidence: str = "",
    predicted_label: str = "local",
    risk_level: str = "low",
    stage_notes: Optional[List[str]] = None,
    core_meta: Optional[Dict[str, Any]] = None,
    remodel_meta: Optional[Dict[str, Any]] = None,
    review_meta: Optional[Dict[str, Any]] = None,
    translation_meta: Optional[Dict[str, Any]] = None,
    timings_ms: Optional[Dict[str, Any]] = None,
    cache_hit: str = "false",
) -> Dict[str, Any]:
    english = str(remodeled_english if remodeled_english is not None else raw_english).strip()
    return {
        "pipeline_version": PIPELINE_VERSION,
        "raw_english": str(raw_english or "").strip(),
        "remodeled_english": english,
        "tamil_text": str(tamil_text or "").strip(),
        "theni_tamil_text": str(theni_tamil_text or tamil_text or "").strip(),
        "direct_answer_source": str(direct_answer_source or ""),
        "direct_answer_confidence": str(direct_answer_confidence or ""),
        "predicted_label": str(predicted_label or "local"),
        "risk_level": str(risk_level or "low"),
        "route_taken": str(route_taken or "local_rag"),
        "cache_hit": str(cache_hit or "false"),
        "stage_notes": json.dumps(stage_notes or [], ensure_ascii=False),
        "core_meta": json.dumps(core_meta or {}, ensure_ascii=False),
        "remodel_meta": json.dumps(remodel_meta or {}, ensure_ascii=False),
        "review_meta": json.dumps(review_meta or {}, ensure_ascii=False),
        "translation_meta": json.dumps(translation_meta or {}, ensure_ascii=False),
        "timings_ms": json.dumps(timings_ms or {"total_ms": 0.0}, ensure_ascii=False),
    }


def _coerce_cached_pipeline(payload: Any) -> Optional[Dict[str, Any]]:
    if not isinstance(payload, dict):
        return None
    pipeline = payload.get("pipeline")
    if not isinstance(pipeline, dict):
        return None

    return _build_pipeline_result(
        raw_english=str(pipeline.get("raw_english", "")).strip(),
        remodeled_english=str(pipeline.get("remodeled_english", "")).strip() or str(pipeline.get("raw_english", "")).strip(),
        tamil_text=str(pipeline.get("tamil_text", "")).strip(),
        theni_tamil_text=str(pipeline.get("theni_tamil_text", "")).strip(),
        route_taken=str(pipeline.get("route_taken", "cached_answer")).strip() or "cached_answer",
        direct_answer_source=str(pipeline.get("direct_answer_source", "qa_cache")).strip() or "qa_cache",
        direct_answer_confidence=str(pipeline.get("direct_answer_confidence", "1.0000")).strip() or "1.0000",
        predicted_label=str(pipeline.get("predicted_label", "cached")).strip() or "cached",
        risk_level=str(pipeline.get("risk_level", "low")).strip() or "low",
        stage_notes=["Reused a cached answer and skipped a new OpenAI call."],
        core_meta=pipeline.get("core_meta") if isinstance(pipeline.get("core_meta"), dict) else {},
        remodel_meta=pipeline.get("remodel_meta") if isinstance(pipeline.get("remodel_meta"), dict) else {},
        review_meta=pipeline.get("review_meta") if isinstance(pipeline.get("review_meta"), dict) else {},
        translation_meta=pipeline.get("translation_meta") if isinstance(pipeline.get("translation_meta"), dict) else {},
        timings_ms=pipeline.get("timings_ms") if isinstance(pipeline.get("timings_ms"), dict) else {"total_ms": 0.0},
    )


def _get_user_timezone(user: Optional[User]) -> ZoneInfo:
    tz_name = (user.timezone if user and user.timezone else "Asia/Kolkata").strip() or "Asia/Kolkata"
    try:
        return ZoneInfo(tz_name)
    except Exception:
        return ZoneInfo("Asia/Kolkata")


def _parse_item_datetime(raw_value: Optional[str], user: Optional[User]) -> Optional[datetime]:
    raw = str(raw_value or "").strip()
    if not raw:
        return None

    try:
        parsed = datetime.fromisoformat(raw.replace("Z", "+00:00"))
    except Exception:
        return None

    if parsed.tzinfo is None:
        return parsed.replace(tzinfo=_get_user_timezone(user))

    return parsed.astimezone(_get_user_timezone(user))


def _format_item_time(item: Item, user: Optional[User]) -> str:
    parsed = _parse_item_datetime(item.datetime_str, user)
    if parsed is None:
        return "Any time"
    return parsed.strftime("%I:%M %p").lstrip("0")


def _collect_items_for_scope(session: Session, user_id: int, scope: str, user: Optional[User]) -> List[Item]:
    all_items = list(session.exec(select(Item).where(Item.user_id == user_id)).all())
    now_local = datetime.now(_get_user_timezone(user))
    today = now_local.date()
    tomorrow = today + timedelta(days=1)
    upcoming_cutoff = now_local + timedelta(days=30)
    collected: List[tuple[datetime, Item]] = []

    for item in all_items:
        parsed = _parse_item_datetime(item.datetime_str, user)
        if parsed is None:
            continue

        include = False
        if scope == "today":
            include = parsed.date() == today
        elif scope == "tomorrow":
            include = parsed.date() == tomorrow
        else:
            include = parsed >= now_local and parsed <= upcoming_cutoff

        if include:
            collected.append((parsed, item))

    collected.sort(key=lambda pair: pair[0])
    return [item for _, item in collected[:5]]


def _build_schedule_answer(session: Session, user_id: int, normalized_query: str, user: Optional[User]) -> Optional[Dict[str, Any]]:
    if not any(word in normalized_query for word in ["schedule", "reminder", "reminders", "task", "tasks", "todo", "plan"]):
        return None

    scope = "upcoming"
    label_en = "upcoming"
    label_ta = "வரவிருக்கும்"

    if "today" in normalized_query or "இன்று" in normalized_query:
        scope = "today"
        label_en = "today"
        label_ta = "இன்று"
    elif "tomorrow" in normalized_query or "நாளை" in normalized_query:
        scope = "tomorrow"
        label_en = "tomorrow"
        label_ta = "நாளை"

    items = _collect_items_for_scope(session, user_id, scope, user)
    display_name = (user.name if user and user.name else "there").strip() or "there"

    if not items:
        english = f"You do not have any {label_en} reminders, {display_name}."
        tamil = f"{display_name}, உங்களுக்கு {label_ta} எந்த நினைவூட்டலும் இல்லை."
        return _build_pipeline_result(
            raw_english=english,
            remodeled_english=english,
            tamil_text=tamil,
            theni_tamil_text=tamil,
            route_taken="local_schedule_rag",
            direct_answer_source="local_schedule_memory",
            direct_answer_confidence="1.0000",
            predicted_label="schedule",
            stage_notes=["Answered from the user's saved reminders without calling OpenAI."],
            timings_ms={"total_ms": 0.0},
        )

    english_lines = [f"You have {len(items)} {label_en} reminder(s), {display_name}:"]
    tamil_lines = [f"{display_name}, உங்களுக்கு {label_ta} {len(items)} நினைவூட்டல்(கள்) இருக்கின்றன:"]

    for idx, item in enumerate(items, start=1):
        title = (item.title or item.raw_text or "Untitled").strip()
        time_label = _format_item_time(item, user)
        english_lines.append(f"{idx}. {time_label} - {title}")
        tamil_lines.append(f"{idx}. {time_label} - {title}")

    english = "\n".join(english_lines)
    tamil = "\n".join(tamil_lines)
    return _build_pipeline_result(
        raw_english=english,
        remodeled_english=english,
        tamil_text=tamil,
        theni_tamil_text=tamil,
        route_taken="local_schedule_rag",
        direct_answer_source="local_schedule_memory",
        direct_answer_confidence="1.0000",
        predicted_label="schedule",
        stage_notes=["Answered from the user's saved reminders without calling OpenAI."],
        timings_ms={"total_ms": 0.0},
    )


def _try_local_fast_path(session: Session, user_id: Optional[int], message: str) -> Optional[Dict[str, Any]]:
    normalized = _normalize_lookup_text(message)
    if not normalized:
        return None

    user = session.get(User, user_id) if user_id else None
    display_name = (user.name if user and user.name else "there").strip() or "there"
    assistant_name = (user.assistant_name if user and user.assistant_name else "Elli").strip() or "Elli"

    greeting_phrases = {"hi", "hey", "hello", "hai", "vanakkam", "வணக்கம்", "ஹலோ"}
    if normalized in greeting_phrases or normalized.startswith("good morning") or normalized.startswith("good evening") or normalized.startswith("good afternoon"):
        if normalized.startswith("good morning"):
            english = f"Good morning {display_name}, how are you doing?"
            tamil = f"காலை வணக்கம் {display_name}, எப்படி இருக்கீங்க?"
        elif normalized.startswith("good evening"):
            english = f"Good evening {display_name}, how are you doing?"
            tamil = f"மாலை வணக்கம் {display_name}, எப்படி இருக்கீங்க?"
        elif normalized.startswith("good afternoon"):
            english = f"Good afternoon {display_name}, how are you doing?"
            tamil = f"மதிய வணக்கம் {display_name}, எப்படி இருக்கீங்க?"
        else:
            english = f"Hi {display_name}, how are you doing?"
            tamil = f"ஹாய் {display_name}, எப்படி இருக்கீங்க?"

        return _build_pipeline_result(
            raw_english=english,
            remodeled_english=english,
            tamil_text=tamil,
            theni_tamil_text=tamil,
            route_taken="local_greeting",
            direct_answer_source="instant_greeting_rule",
            direct_answer_confidence="1.0000",
            predicted_label="greeting",
            stage_notes=["Answered with an instant local greeting without calling OpenAI."],
            timings_ms={"total_ms": 0.0},
        )

    if normalized in {"how are you", "how r you", "epdi iruka", "எப்படி இருக்கீங்க"}:
        english = f"I am doing well, {display_name}. How can I help you today?"
        tamil = f"நான் நல்லா இருக்கேன் {display_name}. இன்று என்ன உதவி வேண்டும்?"
        return _build_pipeline_result(
            raw_english=english,
            remodeled_english=english,
            tamil_text=tamil,
            theni_tamil_text=tamil,
            route_taken="local_smalltalk",
            direct_answer_source="instant_smalltalk_rule",
            direct_answer_confidence="1.0000",
            predicted_label="smalltalk",
            stage_notes=["Answered a small-talk query locally without calling OpenAI."],
            timings_ms={"total_ms": 0.0},
        )

    if normalized in {"thanks", "thank you", "nandri", "நன்றி"}:
        english = f"You're welcome, {display_name}."
        tamil = f"பரவாயில்லை {display_name}, உதவியது சந்தோஷம்."
        return _build_pipeline_result(
            raw_english=english,
            remodeled_english=english,
            tamil_text=tamil,
            theni_tamil_text=tamil,
            route_taken="local_smalltalk",
            direct_answer_source="instant_thanks_rule",
            direct_answer_confidence="1.0000",
            predicted_label="smalltalk",
            stage_notes=["Answered a thank-you query locally without calling OpenAI."],
            timings_ms={"total_ms": 0.0},
        )

    if normalized in {"what is my name", "whats my name", "who am i", "my name"}:
        english = f"Your name is {display_name}."
        tamil = f"உங்கள் பெயர் {display_name}."
        return _build_pipeline_result(
            raw_english=english,
            remodeled_english=english,
            tamil_text=tamil,
            theni_tamil_text=tamil,
            route_taken="local_profile_rag",
            direct_answer_source="local_profile_memory",
            direct_answer_confidence="1.0000",
            predicted_label="profile",
            stage_notes=["Answered from the saved user profile without calling OpenAI."],
            timings_ms={"total_ms": 0.0},
        )

    if normalized in {"who are you", "what is your name", "whats your name", "your name", "what can you do", "help"}:
        english = f"I'm {assistant_name}, your assistant. I can help with reminders, schedules, and quick answers."
        tamil = f"நான் {assistant_name}. நினைவூட்டல்கள், அட்டவணை, மற்றும் விரைவு பதில்களில் நான் உதவ முடியும்."
        return _build_pipeline_result(
            raw_english=english,
            remodeled_english=english,
            tamil_text=tamil,
            theni_tamil_text=tamil,
            route_taken="local_assistant_identity",
            direct_answer_source="local_assistant_profile",
            direct_answer_confidence="1.0000",
            predicted_label="assistant_identity",
            stage_notes=["Answered from app configuration without calling OpenAI."],
            timings_ms={"total_ms": 0.0},
        )

    if user_id:
        schedule_result = _build_schedule_answer(session, user_id, normalized, user)
        if schedule_result is not None:
            return schedule_result

        cache_rows = list(
            session.exec(
                select(QACache).where(QACache.user_id == user_id).order_by(QACache.updated_at.desc())
            ).all()
        )[:40]

        best_payload: Optional[Dict[str, Any]] = None
        best_score = 0.0
        for row in cache_rows:
            score = _token_overlap_score(normalized, row.question)
            if _normalize_lookup_text(row.question) == normalized:
                score = 1.0
            if score > best_score:
                try:
                    payload = json.loads(row.answer or "{}")
                except Exception:
                    payload = {}
                best_payload = _coerce_cached_pipeline(payload)
                best_score = score

        if best_payload and best_score >= 0.96:
            best_payload["route_taken"] = "cached_answer"
            best_payload["direct_answer_source"] = "qa_cache"
            best_payload["direct_answer_confidence"] = f"{best_score:.4f}"
            best_payload["cache_hit"] = "true"
            best_payload["stage_notes"] = json.dumps(
                ["Reused a cached answer and skipped a new OpenAI call."], ensure_ascii=False
            )
            return best_payload

    return None

def _ensure_user_table_auth_columns() -> None:
    inspector = sa_inspect(engine)
    if not inspector.has_table("user"):
        return

    columns = {column["name"] for column in inspector.get_columns("user")}

    with engine.begin() as conn:
        if "firebase_uid" not in columns:
            conn.execute(sa_text('ALTER TABLE "user" ADD COLUMN firebase_uid VARCHAR'))

        if "email" not in columns:
            conn.execute(sa_text('ALTER TABLE "user" ADD COLUMN email VARCHAR'))

        if "reply_language" not in columns:
            conn.execute(sa_text('ALTER TABLE "user" ADD COLUMN reply_language VARCHAR DEFAULT \'ta\''))

        try:
            conn.execute(
                sa_text(
                    'CREATE UNIQUE INDEX IF NOT EXISTS ix_user_firebase_uid_unique ON "user" (firebase_uid)'
                )
            )
        except Exception as exc:
            print(f"[WARN] Could not create firebase_uid index: {exc}")

        try:
            conn.execute(
                sa_text('CREATE UNIQUE INDEX IF NOT EXISTS ix_user_email_unique ON "user" (email)')
            )
        except Exception as exc:
            print(f"[WARN] Could not create email index: {exc}")


@app.on_event("startup")
def ensure_runtime_schema() -> None:
    try:
        _ensure_user_table_auth_columns()
    except Exception as exc:
        print(f"[WARN] Runtime schema sync skipped: {exc}")

@app.middleware("http")
async def log_requests(request: Request, call_next):
    start = datetime.utcnow()
    body_text = ""
    try:
        body = await request.body()
        if body:
            body_text = body.decode("utf-8", errors="ignore")[:1500]
    except Exception:
        pass
    response = await call_next(request)
    ms = int((datetime.utcnow() - start).total_seconds() * 1000)
    print(f"[REQ] {request.method} {request.url.path} {response.status_code} {ms}ms body={body_text[:300]}")
    return response


@app.get("/")
def root():
    return {
        "ok": True,
        "app": "J AI",
        "message": "Persona-aware Tamil assistant backend is running.",
        "endpoints": {
            "health": "/health",
            "text": "/api/chat",
            "voice": "/transcribe-and-analyze",
        },
    }


@app.get("/health")
def health():
    return {"status": "ok", "app": "J AI", "pipeline_version": PIPELINE_VERSION}


@app.get("/api/health")
def api_health():
    return {
        "status": "ok",
        "app": "J AI",
        "mode": PIPELINE_VERSION,
        "features": [
            "persona_context",
            "openai_core_answer",
            "english_remodel",
            "tamil_translation",
            "theni_tamil_conversion",
            "whisper_audio_transcription",
        ],
    }


PARSE_DT_PROMPT = """
You convert natural language datetime into ISO datetime.
Input includes:
- timezone (IANA)
- now (ISO)
- text

Return ONLY JSON:
{
  "iso": "YYYY-MM-DDTHH:MM:SS" or null,
  "human": "human readable summary",
  "confidence": 0.0 to 1.0
}
"""

SYSTEM_PROMPT = """
You are a personal AI assistant.

You will receive JSON input with:
- context:
    - user (name, place, timezone)
    - routine (daily schedule and habits)
    - personality (communication style, motivation, sensitivity)
- input (the user's message)

Rules:
- Always respect the user's routine when suggesting times or actions
- Always match your tone to the personality profile
- If personality is missing, be neutral and polite
- If routine is missing, ask clarifying questions
- Never suggest actions outside wake/sleep boundaries unless explicitly asked
- If a request conflicts with routine, explain and suggest an alternative

Return ONLY JSON:
{
  "intent": "reminder|note|task|document|other",
  "category": "Work|Home|Business|Other",
  "datetime": "... or null",
  "title": "...",
  "details": "..."
}
"""

PERSONALITY_QUESTIONS = [
    "How would you describe yourself in one sentence?",
    "Do you prefer strict reminders or gentle nudges?",
    "Are you more spontaneous or planned?",
    "What usually motivates you?",
    "How do you want the assistant to talk to you?",
]

PERSONALITY_SUMMARY_PROMPT = """
You are analyzing a user's personality.

Based on their answers, create a concise profile including:
- communication tone
- motivation style
- structure vs flexibility preference
- emotional sensitivity

Return plain text. No JSON.
"""

CHECKIN_PROMPT = """
- Match message tone to personality
- If personality prefers gentle nudges, avoid commands
- If personality prefers strictness, be direct

You will receive:
- user profile (name, place, timezone)
- user routine (wake_time, sleep_time, work_start, work_end, daily_habits)

Create 3–8 smart check-ins for today.
Respect the user's routine and time boundaries.

Each check-in must include:
- title
- when (HH:MM 24h)
- message (address user by name)

Return ONLY JSON:
{ "checkins": [ { "title":"...", "when":"08:00", "message":"..." } ] }
"""


class ParseDatetimeRequest(BaseModel):
    text: str
    timezone: str = "Asia/Kolkata"
    now_iso: Optional[str] = None


class DailyRoutineIn(BaseModel):
    wake_time: str
    sleep_time: str
    work_start: Optional[str] = None
    work_end: Optional[str] = None
    daily_habits: Optional[str] = None


class DailyRoutineOut(DailyRoutineIn):
    user_id: int


class TextAnalysisRequest(BaseModel):
    text: str
    user_id: Optional[int] = None
    meta: Optional[Dict[str, Any]] = None
    reply_language: Optional[str] = None


class PersonalityAnswersIn(BaseModel):
    answers: Dict[str, str]


class TextAnalysisResponse(BaseModel):
    id: int
    intent: str
    category: str
    raw_text: str
    transcript: Optional[str] = None
    datetime: Optional[str] = None
    title: Optional[str] = None
    details: Optional[str] = None


class UserCreate(BaseModel):
    user_id: Optional[int] = None
    firebase_uid: Optional[str] = None
    email: Optional[str] = None
    name: str
    place: Optional[str] = None
    timezone: Optional[str] = "Asia/Kolkata"
    assistant_name: Optional[str] = "Elli"
    reply_language: Optional[str] = "ta"


class ChatAPIRequest(BaseModel):
    user_id: Optional[int] = None
    message: Optional[str] = None
    text: Optional[str] = None
    include_pipeline: bool = True
    reply_language: Optional[str] = None


class PipelineChatRequest(BaseModel):
    user_id: int
    message: str


def _extract_response_text(response: Any) -> str:
    output_text = getattr(response, "output_text", None)
    if output_text:
        return str(output_text).strip()

    chunks: List[str] = []
    for item in getattr(response, "output", None) or []:
        for part in getattr(item, "content", None) or []:
            text = getattr(part, "text", None)
            if text:
                chunks.append(str(text))
            elif isinstance(part, dict) and part.get("text"):
                chunks.append(str(part["text"]))
    return "\n".join(part.strip() for part in chunks if str(part).strip()).strip()


def llm_json(system_prompt: str, user_content: str, temperature: float = 0.2) -> Dict[str, Any]:
    response = client.responses.create(
        model=OPENAI_JSON_MODEL,
        input=[
            {"role": "system", "content": [{"type": "input_text", "text": system_prompt.strip()}]},
            {"role": "user", "content": [{"type": "input_text", "text": user_content.strip()}]},
        ],
        temperature=temperature,
        text={"format": {"type": "json_object"}},
    )
    raw = _extract_response_text(response)
    return json.loads(raw)


def llm_text(system_prompt: str, user_content: str, temperature: float = 0.2) -> str:
    response = client.responses.create(
        model=OPENAI_JSON_MODEL,
        input=[
            {"role": "system", "content": [{"type": "input_text", "text": system_prompt.strip()}]},
            {"role": "user", "content": [{"type": "input_text", "text": user_content.strip()}]},
        ],
        temperature=temperature,
    )
    return _extract_response_text(response)


def normalize_category(raw: str) -> str:
    cr = (raw or "Other").lower()
    if cr == "work":
        return "Work"
    if cr == "home":
        return "Home"
    if cr == "business":
        return "Business"
    return "Other"


def validate_hhmm(v: str) -> None:
    if not re.match(r"^([01]\d|2[0-3]):[0-5]\d$", str(v or "").strip()):
        raise HTTPException(400, f"Invalid time format: {v}")


def normalize_optional(v: Optional[str]) -> Optional[str]:
    if v is None:
        return None
    v = str(v).strip()
    return v if v else None


def item_to_response(item: Item) -> TextAnalysisResponse:
    return TextAnalysisResponse(
        id=item.id,
        intent=item.intent,
        category=item.category,
        raw_text=item.raw_text,
        transcript=item.transcript,
        datetime=item.datetime_str,
        title=item.title,
        details=item.details,
    )


def log_conversation(
    session: Session,
    user_id: Optional[int],
    channel: str,
    user_input: str,
    transcript: Optional[str],
    llm_json_out: Optional[dict],
):
    row = Conversation(
        user_id=user_id,
        channel=channel,
        user_input=user_input,
        transcript=transcript,
        llm_output_json=json.dumps(llm_json_out, ensure_ascii=False) if llm_json_out else None,
        created_at=datetime.utcnow(),
    )
    session.add(row)
    session.commit()


def upsert_qa_cache(session: Session, user_id: Optional[int], question: str, answer_json: dict):
    q = select(QACache).where(QACache.question == question)
    if user_id is not None:
        q = q.where(QACache.user_id == user_id)
    row = session.exec(q).first()
    if row:
        row.answer = json.dumps(answer_json, ensure_ascii=False)
        row.hits = (row.hits or 0) + 1
        row.updated_at = datetime.utcnow()
        session.add(row)
        session.commit()
        return
    session.add(
        QACache(
            user_id=user_id,
            question=question,
            answer=json.dumps(answer_json, ensure_ascii=False),
            hits=1,
            updated_at=datetime.utcnow(),
        )
    )
    session.commit()


def build_user_context(session: Session, user_id: int) -> dict:
    user = session.get(User, user_id)
    if not user:
        return {}

    routine = session.exec(select(DailyRoutine).where(DailyRoutine.user_id == user_id)).first()
    profile = session.exec(select(UserProfile).where(UserProfile.user_id == user_id)).first()

    if profile and profile.questions_version != PERSONALITY_QUESTIONS_VERSION:
        personality = "Personality profile outdated. Be neutral and helpful."
    elif profile and profile.profile_summary:
        personality = profile.profile_summary
    else:
        personality = "No personality profile yet. Be neutral and helpful."

    return {
        "user": {
            "name": user.name,
            "place": user.place,
            "timezone": user.timezone,
        },
        "routine": {
            "wake_time": routine.wake_time if routine else None,
            "sleep_time": routine.sleep_time if routine else None,
            "work_start": routine.work_start if routine else None,
            "work_end": routine.work_end if routine else None,
            "daily_habits": routine.daily_habits if routine else None,
        },
        "personality": personality,
    }


def _infer_hobbies(habits_text: str) -> List[str]:
    s = (habits_text or "").lower()
    hobbies = []
    if "walk" in s:
        hobbies.append("walking")
    if "read" in s:
        hobbies.append("reading")
    if "gym" in s or "workout" in s:
        hobbies.append("fitness")
    if "pray" in s:
        hobbies.append("spirituality")
    if "cook" in s:
        hobbies.append("cooking")
    if "music" in s:
        hobbies.append("music")
    return hobbies[:4]


def _infer_tone(summary: str) -> str:
    s = (summary or "").lower()
    if any(word in s for word in ["short", "brief"]):
        return "brief"
    if any(word in s for word in ["detail", "deep"]):
        return "detailed"
    if any(word in s for word in ["casual", "friendly"]):
        return "friendly_casual"
    if any(word in s for word in ["formal", "respectful"]):
        return "respectful"
    return "warm"


def _infer_personality_style(summary: str) -> str:
    s = (summary or "").lower()
    if "calm" in s:
        return "calm"
    if "friendly" in s:
        return "friendly"
    if "ambitious" in s or "goal" in s:
        return "ambitious"
    if "sensitive" in s or "emotional" in s:
        return "emotional_sensitive"
    return "practical"


def _default_stage_answers(
    user: Optional[User],
    routine: Optional[DailyRoutine],
    db_profile: Optional[UserProfile],
) -> Dict[str, Any]:
    habits_text = routine.daily_habits if routine else ""
    summary = db_profile.profile_summary if db_profile and db_profile.profile_summary else ""
    return {
        "age_group": "26-35",
        "gender_context": "prefer_not_to_say",
        "life_stage": "none_of_these",
        "food_preference": "mixed_flexible",
        "health_conditions": ["none"],
        "food_caution": "no_special_caution",
        "daily_activity": "moderate_walks",
        "sleep_pattern": "average",
        "personality_style": _infer_personality_style(summary),
        "stress_support": "step_by_step_plan",
        "communication_tone": _infer_tone(summary),
        "answer_length": "medium",
        "hobbies": _infer_hobbies(habits_text or ""),
        "main_goal": "career_or_business",
        "family_role": "working_professional",
    }


def _sync_stage_profile(session: Session, user_id: Optional[int]) -> Dict[str, Any]:
    uid = str(user_id or "guest")
    user = session.get(User, user_id) if user_id else None
    routine = session.exec(select(DailyRoutine).where(DailyRoutine.user_id == user_id)).first() if user_id else None
    db_profile = session.exec(select(UserProfile).where(UserProfile.user_id == user_id)).first() if user_id else None

    answers = _default_stage_answers(user, routine, db_profile)
    profile = {
        "profile_version": "ai_tool_db_bridge",
        "user_id": uid,
        "created_at": datetime.utcnow().isoformat() + "Z",
        "answers": answers,
        "behaviour_rules": STAGE_BEHAVIOUR._derive_behaviour_rules(answers),
        "rag_personality_hints": STAGE_BEHAVIOUR._infer_personality_rag_from_answers(answers),
        "app_profile": {
            "name": user.name if user else "Guest",
            "place": user.place if user else "",
            "timezone": user.timezone if user else "Asia/Kolkata",
            "assistant_name": user.assistant_name if user else "Elli",
            "reply_language": _normalize_reply_language(user.reply_language if user else "ta"),
        },
        "daily_routine": {
            "wake_time": routine.wake_time if routine else "07:30",
            "sleep_time": routine.sleep_time if routine else "23:30",
            "work_start": routine.work_start if routine else "09:30",
            "work_end": routine.work_end if routine else "18:30",
            "daily_habits": routine.daily_habits if routine else "",
        },
    }

    existing = None
    if STAGE_BEHAVIOUR.profile_exists(uid):
        try:
            existing = STAGE_BEHAVIOUR.load_profile(uid)
        except Exception:
            existing = None
    if existing:
        profile["created_at"] = existing.get("created_at", profile["created_at"])

    profile = STAGE_BEHAVIOUR._upgrade_profile(profile)
    if db_profile and db_profile.profile_summary:
        stage_summary = str(profile.get("profile_summary", "")).strip()
        profile["profile_summary"] = (
            f"{stage_summary}\n\nExisting app personality summary: {db_profile.profile_summary.strip()}".strip()
        )

    STAGE_BEHAVIOUR.save_profile(uid, profile)
    return profile


def _safe_json(value: Any) -> str:
    try:
        return json.dumps(value, ensure_ascii=False)
    except Exception:
        return json.dumps(str(value), ensure_ascii=False)


def _log_stage_history(user_id: Optional[int], profile: Dict[str, Any], query: str, result: Dict[str, Any]) -> None:
    uid = str(user_id or "guest")
    log_path = LOGS_DIR / f"{uid}_history.jsonl"
    record = {
        "timestamp": datetime.utcnow().isoformat() + "Z",
        "user_id": uid,
        "query": query,
        "profile_summary": profile.get("profile_summary", ""),
        "profile_card": profile.get("profile_card", {}),
        "result": result,
    }
    with log_path.open("a", encoding="utf-8") as handle:
        handle.write(json.dumps(record, ensure_ascii=False) + "\n")


def _run_stage_pipeline(session: Session, user_id: Optional[int], message: str, reply_language: Optional[str] = None) -> Dict[str, Any]:
    pipeline_user = session.get(User, user_id) if user_id else None
    resolved_reply_language = _normalize_reply_language(reply_language or (pipeline_user.reply_language if pipeline_user else None))
    uid = str(user_id or "guest")
    cache_key = f"{uid}:{resolved_reply_language}::{' '.join(message.strip().lower().split())}"
    if cache_key in STAGE_CACHE:
        cached = dict(STAGE_CACHE[cache_key])
        cached["cache_hit"] = "true"
        return cached

    fast_path = _try_local_fast_path(session, user_id, message)
    if fast_path is not None:
        STAGE_CACHE[cache_key] = dict(fast_path)
        if len(STAGE_CACHE) > 128:
            first_key = next(iter(STAGE_CACHE))
            STAGE_CACHE.pop(first_key, None)
        return fast_path

    profile = _sync_stage_profile(session, user_id)
    total_start = time.perf_counter()
    timings: Dict[str, float] = {}
    stage_notes: List[str] = []

    t0 = time.perf_counter()
    profile_context = STAGE_BEHAVIOUR.build_runtime_context(profile, user_query=message)
    timings["context_ms"] = round((time.perf_counter() - t0) * 1000, 2)

    t0 = time.perf_counter()
    direct_match = STAGE_REMODELER.get_direct_answer_match(message)
    timings["direct_match_ms"] = round((time.perf_counter() - t0) * 1000, 2)

    core_meta: Dict[str, Any] = {"answer": "", "answer_style": "", "risk_level": "low", "safety_notes": ""}
    remodel_meta: Dict[str, Any] = {}
    review_meta: Dict[str, Any] = {}
    translation_meta: Dict[str, Any] = {}
    route_taken = "full_pipeline"
    direct_answer_source = ""
    direct_answer_confidence = ""
    predicted_label = "unknown"
    risk_level = "low"

    if direct_match and direct_match.confidence >= 0.92:
        raw_english = direct_match.answer
        remodeled_english = direct_match.answer
        route_taken = "dataset_direct_answer"
        direct_answer_source = f"{direct_match.match_type}:{direct_match.query}"
        direct_answer_confidence = f"{direct_match.confidence:.4f}"
        predicted_label = direct_match.label
        stage_notes.append("Used a high-confidence direct answer from the local dataset.")
    else:
        t0 = time.perf_counter()
        core_meta = STAGE_CORE.answer_user_query_structured(message, profile_context)
        timings["core_answer_ms"] = round((time.perf_counter() - t0) * 1000, 2)
        raw_english = str(core_meta.get("answer", "")).strip()

        t0 = time.perf_counter()
        remodel_meta = STAGE_REMODELER.remodel_with_meta(message, raw_english, profile)
        timings["remodel_ms"] = round((time.perf_counter() - t0) * 1000, 2)
        remodeled_english = str(remodel_meta.get("answer", raw_english)).strip() or raw_english

        t0 = time.perf_counter()
        review_meta = STAGE_CORE.review_answer(message, remodeled_english, profile_context)
        timings["review_ms"] = round((time.perf_counter() - t0) * 1000, 2)
        remodeled_english = str(review_meta.get("final_answer", remodeled_english)).strip() or remodeled_english

        route_taken = str(remodel_meta.get("route", "full_rewrite"))
        predicted_label = str(remodel_meta.get("predicted_label", "unknown"))
        risk_level = str(remodel_meta.get("risk_level") or core_meta.get("risk_level") or "low")
        direct_answer_source = str(remodel_meta.get("direct_answer_source", ""))
        if remodel_meta.get("direct_answer_confidence") not in (None, ""):
            direct_answer_confidence = f"{float(remodel_meta.get('direct_answer_confidence', 0.0)):.4f}"

        for note in (core_meta.get("safety_notes"), remodel_meta.get("route_reason"), review_meta.get("review_note")):
            if str(note or "").strip():
                stage_notes.append(str(note).strip())

    tamil_text = ""
    theni_tamil_text = ""
    if resolved_reply_language == "ta":
        t0 = time.perf_counter()
        translation_meta = STAGE_TRANSLATOR.english_to_tamil_with_meta(remodeled_english, profile)
        tamil_text = str(translation_meta.get("tamil_text", "")).strip()
        timings["english_to_tamil_ms"] = round((time.perf_counter() - t0) * 1000, 2)

        t0 = time.perf_counter()
        theni_tamil_text = STAGE_TRANSLATOR.tamil_to_thenitamil(tamil_text)
        timings["tamil_to_theni_ms"] = round((time.perf_counter() - t0) * 1000, 2)
    else:
        translation_meta = {"skipped": True, "reason": "reply_language_is_english"}

    total_ms = round((time.perf_counter() - total_start) * 1000, 2)
    result: Dict[str, Any] = {
        "pipeline_version": PIPELINE_VERSION,
        "raw_english": raw_english,
        "remodeled_english": remodeled_english,
        "tamil_text": tamil_text,
        "theni_tamil_text": theni_tamil_text,
        "direct_answer_source": direct_answer_source,
        "direct_answer_confidence": direct_answer_confidence,
        "predicted_label": predicted_label,
        "risk_level": risk_level,
        "route_taken": route_taken,
        "cache_hit": "false",
        "stage_notes": _safe_json(stage_notes),
        "core_meta": _safe_json(core_meta),
        "remodel_meta": _safe_json(remodel_meta),
        "review_meta": _safe_json(review_meta),
        "translation_meta": _safe_json(translation_meta),
        "timings_ms": json.dumps({**timings, "total_ms": total_ms}, ensure_ascii=False),
    }

    _log_stage_history(user_id, profile, message, result)
    STAGE_CACHE[cache_key] = dict(result)
    if len(STAGE_CACHE) > 128:
        first_key = next(iter(STAGE_CACHE))
        STAGE_CACHE.pop(first_key, None)
    return result


def _metadata_for_item(session: Session, user_id: Optional[int], text: str, fallback_details: str) -> Dict[str, Any]:
    user_context = {}
    if user_id:
        user_context = build_user_context(session, user_id)
    if user_id and not user_context.get("personality"):
        user_context["personality"] = "Unknown personality. Be neutral and helpful."

    user_content = json.dumps({"context": user_context, "input": text}, ensure_ascii=False)

    try:
        data = llm_json(SYSTEM_PROMPT, user_content, temperature=0.2)
        if not data.get("details"):
            data["details"] = fallback_details
        return data
    except Exception:
        clean_text = " ".join(text.strip().split())
        return {
            "intent": "other",
            "category": "Other",
            "datetime": None,
            "title": (clean_text[:60] + "...") if len(clean_text) > 60 else clean_text,
            "details": fallback_details,
        }


def _normalized_pipeline_result(result: Dict[str, Any]) -> Dict[str, Any]:
    def _maybe(value: Any, default: Any):
        if isinstance(value, str):
            text = value.strip()
            if text and text[:1] in "[{":
                try:
                    return json.loads(text)
                except Exception:
                    return default
        return value if value not in (None, "") else default

    return {
        "pipeline_version": result.get("pipeline_version", PIPELINE_VERSION),
        "raw_english": result.get("raw_english", ""),
        "remodeled_english": result.get("remodeled_english", ""),
        "tamil_text": result.get("tamil_text", ""),
        "theni_tamil_text": result.get("theni_tamil_text", ""),
        "direct_answer_source": result.get("direct_answer_source", ""),
        "direct_answer_confidence": result.get("direct_answer_confidence", ""),
        "predicted_label": result.get("predicted_label", ""),
        "risk_level": result.get("risk_level", ""),
        "route_taken": result.get("route_taken", ""),
        "cache_hit": result.get("cache_hit", "false"),
        "stage_notes": _maybe(result.get("stage_notes"), []),
        "core_meta": _maybe(result.get("core_meta"), {}),
        "remodel_meta": _maybe(result.get("remodel_meta"), {}),
        "review_meta": _maybe(result.get("review_meta"), {}),
        "translation_meta": _maybe(result.get("translation_meta"), {}),
        "timings_ms": _maybe(result.get("timings_ms"), {}),
    }


def _assistant_text_from_pipeline(
    pipeline_result: Dict[str, Any],
    fallback: str,
    reply_language: Optional[str] = None,
) -> str:
    resolved_reply_language = _normalize_reply_language(reply_language)
    if resolved_reply_language == "en":
        return str(pipeline_result.get("remodeled_english") or "").strip() or fallback

    return (
        str(pipeline_result.get("theni_tamil_text") or "").strip()
        or str(pipeline_result.get("tamil_text") or "").strip()
        or str(pipeline_result.get("remodeled_english") or "").strip()
        or fallback
    )


def _save_item_from_pipeline(
    session: Session,
    *,
    user_id: Optional[int],
    source: str,
    raw_text: str,
    transcript: Optional[str],
    pipeline_result: Dict[str, Any],
    reply_language: Optional[str] = None,
) -> tuple[Item, Dict[str, Any], Dict[str, Any]]:
    spoken_answer = _assistant_text_from_pipeline(pipeline_result, raw_text, reply_language)
    meta = _metadata_for_item(session, user_id, raw_text, spoken_answer)

    item = Item(
        intent=str(meta.get("intent", "other")).lower(),
        category=normalize_category(str(meta.get("category", "Other"))),
        raw_text=raw_text,
        transcript=transcript,
        datetime_str=meta.get("datetime"),
        title=meta.get("title") or raw_text[:60],
        details=spoken_answer,
        source=source,
        user_id=user_id,
        created_at=datetime.utcnow(),
        updated_at=datetime.utcnow(),
    )
    session.add(item)
    session.commit()
    session.refresh(item)

    normalized_pipeline = _normalized_pipeline_result(pipeline_result)
    payload = {"pipeline": normalized_pipeline, "meta": meta}
    log_conversation(session, user_id, source, raw_text, transcript, payload)
    upsert_qa_cache(session, user_id, raw_text, payload)
    return item, meta, normalized_pipeline


def _build_chat_response(item: Item, meta: Dict[str, Any], pipeline: Dict[str, Any]) -> Dict[str, Any]:
    assistant_text = item.details or item.raw_text
    return {
        "ok": True,
        "item": {
            "id": item.id,
            "intent": item.intent,
            "category": item.category,
            "title": item.title,
            "details": assistant_text,
            "datetime": item.datetime_str,
            "source": item.source,
            "raw_text": item.raw_text,
            "transcript": item.transcript,
            "created_at": item.created_at.isoformat() if item.created_at else None,
        },
        "assistant": {
            "text": assistant_text,
            "english": pipeline.get("remodeled_english", ""),
            "tamil": pipeline.get("tamil_text", ""),
            "theni_tamil": pipeline.get("theni_tamil_text", ""),
        },
        "pipeline": pipeline,
        "meta": meta,
    }


def _resolve_chat_text(payload: ChatAPIRequest) -> str:
    text = (payload.message or payload.text or "").strip()
    if not text:
        raise HTTPException(400, "message or text is required")
    return text


def _transcribe_audio_file(file_path: str) -> str:
    with open(file_path, "rb") as audio_file:
        transcript_obj = client.audio.transcriptions.create(
            model="whisper-1",
            file=audio_file,
            response_format="json",
            language="ta",
        )
    text = str(getattr(transcript_obj, "text", "") or "").strip()
    if not text:
        raise HTTPException(400, "Failed to transcribe audio")
    return text


@app.post("/parse-datetime")
def parse_datetime(payload: ParseDatetimeRequest):
    now_iso = payload.now_iso or datetime.utcnow().isoformat()
    user_content = json.dumps({"timezone": payload.timezone, "now": now_iso, "text": payload.text}, ensure_ascii=False)
    out = llm_json(PARSE_DT_PROMPT, user_content, temperature=0.0)
    return {
        "iso": out.get("iso"),
        "human": out.get("human") or "",
        "confidence": float(out.get("confidence") or 0.0),
    }
def _normalize_email(email: Optional[str]) -> Optional[str]:
    normalized = (email or "").strip().lower()
    return normalized or None


def _questionnaire_completed(profile: Optional[UserProfile]) -> bool:
    if not profile:
        return False

    try:
        answers = json.loads(profile.answers_json or "{}")
    except Exception:
        answers = {}

    return bool(answers)


def _ensure_user_profile(session: Session, user_id: int) -> UserProfile:
    profile = session.exec(select(UserProfile).where(UserProfile.user_id == user_id)).first()
    if profile:
        return profile

    profile = UserProfile(
        user_id=user_id,
        answers_json=json.dumps({}, ensure_ascii=False),
        questions_version=PERSONALITY_QUESTIONS_VERSION,
        updated_at=datetime.utcnow(),
    )
    session.add(profile)
    session.commit()
    session.refresh(profile)
    return profile


def _serialize_user_payload(user: User, profile: Optional[UserProfile]) -> Dict[str, Any]:
    assistant_name = (user.assistant_name or "Elli").strip() or "Elli"

    return {
        "id": int(user.id) if user.id is not None else None,
        "firebase_uid": user.firebase_uid,
        "email": user.email or "",
        "name": user.name,
        "place": user.place or "",
        "timezone": user.timezone or "Asia/Kolkata",
        "assistant_name": assistant_name,
        "reply_language": _normalize_reply_language(getattr(user, "reply_language", "ta")),
        "created_at": user.created_at.isoformat() if user.created_at else None,
        "profile_id": int(profile.id) if profile and profile.id is not None else None,
        "profile_user_id": int(profile.user_id) if profile and profile.user_id is not None else None,
        "questionnaire_completed": _questionnaire_completed(profile),
    }

@app.post("/users")
def create_user(payload: UserCreate, session: Session = Depends(get_session)):
    assistant_name = (payload.assistant_name or "Elli").strip() or "Elli"
    reply_language = _normalize_reply_language(payload.reply_language)
    normalized_email = _normalize_email(payload.email)
    firebase_uid = (payload.firebase_uid or "").strip() or None

    print(
        "[DEBUG] /users incoming payload:",
        json.dumps(payload.model_dump(), ensure_ascii=False, default=str),
    )

    user: Optional[User] = None

    try:
        if payload.user_id:
            user = session.get(User, payload.user_id)

        if not user and firebase_uid:
            user = session.exec(select(User).where(User.firebase_uid == firebase_uid)).first()

        if not user and normalized_email:
            user = session.exec(select(User).where(User.email == normalized_email)).first()

        if user:
            user.firebase_uid = firebase_uid or user.firebase_uid
            user.email = normalized_email or user.email
            user.name = payload.name
            user.place = payload.place
            user.timezone = payload.timezone or user.timezone or "Asia/Kolkata"
            user.assistant_name = assistant_name
            user.reply_language = reply_language or getattr(user, "reply_language", "ta") or "ta"
            session.add(user)
            session.commit()
            session.refresh(user)
        else:
            user = User(
                firebase_uid=firebase_uid,
                email=normalized_email,
                name=payload.name,
                place=payload.place,
                timezone=payload.timezone or "Asia/Kolkata",
                assistant_name=assistant_name,
                reply_language=reply_language,
            )
            session.add(user)
            session.commit()
            session.refresh(user)

        profile = _ensure_user_profile(session, int(user.id))

    except Exception as exc:
        session.rollback()
        print(f"[DEBUG] /users failed while saving: {exc}")
        raise

    response_payload = _serialize_user_payload(user, profile)

    print(
        "[DEBUG] /users response payload:",
        json.dumps(response_payload, ensure_ascii=False, default=str),
    )

    return response_payload


@app.get("/users/resolve")
def resolve_user(
    firebase_uid: Optional[str] = Query(default=None),
    email: Optional[str] = Query(default=None),
    session: Session = Depends(get_session),
):
    normalized_email = _normalize_email(email)
    normalized_uid = (firebase_uid or "").strip() or None

    if not normalized_uid and not normalized_email:
        raise HTTPException(400, "firebase_uid or email is required")

    user: Optional[User] = None

    if normalized_uid:
        user = session.exec(select(User).where(User.firebase_uid == normalized_uid)).first()

    if not user and normalized_email:
        user = session.exec(select(User).where(User.email == normalized_email)).first()

    if not user:
        return {"found": False}

    profile = _ensure_user_profile(session, int(user.id))

    if normalized_uid and user.firebase_uid != normalized_uid:
        user.firebase_uid = normalized_uid
        session.add(user)
        session.commit()
        session.refresh(user)

    if normalized_email and user.email != normalized_email:
        user.email = normalized_email
        session.add(user)
        session.commit()
        session.refresh(user)

    return {
        "found": True,
        "user": _serialize_user_payload(user, profile),
    }


@app.get("/users/{user_id}")
def get_user(user_id: int, session: Session = Depends(get_session)):
    user = session.get(User, user_id)
    if not user:
        raise HTTPException(404, "User not found")
    profile = _ensure_user_profile(session, user_id)
    return _serialize_user_payload(user, profile)

@app.delete("/users/{user_id}")
def delete_user_account(user_id: int, session: Session = Depends(get_session)):
    user = session.get(User, user_id)
    if not user:
        raise HTTPException(404, "User not found")

    try:
        session.exec(delete(Item).where(Item.user_id == user_id))
        session.exec(delete(Conversation).where(Conversation.user_id == user_id))
        session.exec(delete(QACache).where(QACache.user_id == user_id))
        session.exec(delete(DailyRoutine).where(DailyRoutine.user_id == user_id))
        session.exec(delete(UserProfile).where(UserProfile.user_id == user_id))
        session.delete(user)
        session.commit()
    except Exception as exc:
        session.rollback()
        print(f"[DEBUG] /users/{{user_id}} delete failed: {exc}")
        raise

    for path_getter in (STAGE_BEHAVIOUR._profile_path, STAGE_BEHAVIOUR._history_log_path):
        try:
            path_getter(str(user_id)).unlink(missing_ok=True)
        except Exception as exc:
            print(f"[WARN] Failed to delete stage file for user {user_id}: {exc}")

    keys_to_delete = [key for key in STAGE_CACHE if key.startswith(f"{user_id}::")]
    for key in keys_to_delete:
        STAGE_CACHE.pop(key, None)

    return {"ok": True, "deleted_user_id": user_id}

@app.get("/personality/questions")
def get_personality_questions():
    return {"version": PERSONALITY_QUESTIONS_VERSION, "questions": PERSONALITY_QUESTIONS}


@app.get("/api/questions")
def get_pipeline_questions():
    return {"questions": PIPELINE_QUESTIONS}


@app.get("/api/profile/{user_id}")
def get_pipeline_profile(user_id: int, session: Session = Depends(get_session)):
    profile = _sync_stage_profile(session, user_id)
    return {"exists": True, "profile": profile}


@app.post("/api/profile/{user_id}")
def save_pipeline_profile(user_id: int, payload: PersonalityAnswersIn, session: Session = Depends(get_session)):
    answers_json = json.dumps(payload.answers, ensure_ascii=False)
    profile = session.exec(select(UserProfile).where(UserProfile.user_id == user_id)).first()
    if profile:
        profile.answers_json = answers_json
        profile.profile_summary = None
        profile.questions_version = PERSONALITY_QUESTIONS_VERSION
        profile.updated_at = datetime.utcnow()
    else:
        profile = UserProfile(
            user_id=user_id,
            answers_json=answers_json,
            questions_version=PERSONALITY_QUESTIONS_VERSION,
            updated_at=datetime.utcnow(),
        )
        session.add(profile)
    session.commit()
    session.refresh(profile)
    stage_profile = _sync_stage_profile(session, user_id)
    STAGE_CACHE.clear()
    return {"ok": True, "profile": stage_profile}


@app.post("/api/chat")
def api_chat(payload: ChatAPIRequest, session: Session = Depends(get_session)):
    text = _resolve_chat_text(payload)
    pipeline_result = _run_stage_pipeline(session, payload.user_id, text, payload.reply_language)
    item, meta, normalized_pipeline = _save_item_from_pipeline(
        session,
        user_id=payload.user_id,
        source="text",
        raw_text=text,
        transcript=None,
        pipeline_result=pipeline_result,
        reply_language=payload.reply_language,
    )
    return _build_chat_response(item, meta, normalized_pipeline)


@app.post("/users/{user_id}/questionnaire")
def save_mobile_questionnaire(user_id: int, payload: Dict[str, Any], session: Session = Depends(get_session)):
    raw = payload.get("payload", payload)
    mapped = DailyRoutineIn(
        wake_time=str(raw.get("wake") or raw.get("wake_time") or "07:30"),
        sleep_time=str(raw.get("sleep") or raw.get("sleep_time") or "23:30"),
        work_start=raw.get("workStart") or raw.get("work_start"),
        work_end=raw.get("workEnd") or raw.get("work_end"),
        daily_habits=raw.get("dailyHabits") or raw.get("daily_habits"),
    )
    return upsert_daily_routine(user_id, mapped, session)


@app.post("/users/{user_id}/generate-daily-checkins")
def generate_daily_checkins(user_id: int, session: Session = Depends(get_session)):
    user = session.get(User, user_id)
    if not user:
        raise HTTPException(404, "User not found")
    routine = session.exec(select(DailyRoutine).where(DailyRoutine.user_id == user_id)).first()
    if not routine:
        raise HTTPException(400, "Daily routine not set. Please configure routine first.")
    profile = session.exec(select(UserProfile).where(UserProfile.user_id == user_id)).first()
    user_content = json.dumps(
        {
            "user": {"name": user.name, "place": user.place, "timezone": user.timezone},
            "personality": profile.profile_summary if profile and profile.profile_summary else "No personality profile yet. Be neutral and helpful.",
            "routine": {
                "wake_time": routine.wake_time,
                "sleep_time": routine.sleep_time,
                "work_start": routine.work_start,
                "work_end": routine.work_end,
                "daily_habits": routine.daily_habits,
            },
            "today": str(date.today()),
        },
        ensure_ascii=False,
    )
    out = llm_json(CHECKIN_PROMPT, user_content, temperature=0.2)
    checkins = out.get("checkins", [])
    checkins.sort(key=lambda x: x.get("when", "99:99"))
    log_conversation(session, user_id, "system", "generate-daily-checkins", None, out)
    return {"checkins": checkins}


@app.get("/users/{user_id}/daily-routine", response_model=DailyRoutineOut)
def get_daily_routine(user_id: int, session: Session = Depends(get_session)):
    routine = session.exec(select(DailyRoutine).where(DailyRoutine.user_id == user_id)).first()
    if not routine:
        raise HTTPException(404, "Daily routine not set")
    return routine


@app.put("/users/{user_id}/daily-routine", response_model=DailyRoutineOut)
def upsert_daily_routine(user_id: int, payload: DailyRoutineIn, session: Session = Depends(get_session)):
    work_start = normalize_optional(payload.work_start)
    work_end = normalize_optional(payload.work_end)
    daily_habits = normalize_optional(payload.daily_habits)
    validate_hhmm(payload.wake_time)
    validate_hhmm(payload.sleep_time)
    if work_start:
        validate_hhmm(work_start)
    if work_end:
        validate_hhmm(work_end)

    routine = session.exec(select(DailyRoutine).where(DailyRoutine.user_id == user_id)).first()
    if routine:
        routine.wake_time = payload.wake_time
        routine.sleep_time = payload.sleep_time
        routine.work_start = work_start
        routine.work_end = work_end
        routine.daily_habits = daily_habits
        routine.updated_at = datetime.utcnow()
    else:
        routine = DailyRoutine(
            user_id=user_id,
            wake_time=payload.wake_time,
            sleep_time=payload.sleep_time,
            work_start=work_start,
            work_end=work_end,
            daily_habits=daily_habits,
            updated_at=datetime.utcnow(),
        )
        session.add(routine)
    session.commit()
    session.refresh(routine)
    STAGE_CACHE.clear()
    _sync_stage_profile(session, user_id)
    return routine


@app.get("/users/{user_id}/personality")
def get_personality(user_id: int, session: Session = Depends(get_session)):
    profile = session.exec(select(UserProfile).where(UserProfile.user_id == user_id)).first()
    if not profile:
        raise HTTPException(404, "Personality profile not found")
    return {"answers": json.loads(profile.answers_json or "{}"), "summary": profile.profile_summary}


@app.post("/users/{user_id}/personality")
def save_personality_answers(user_id: int, payload: PersonalityAnswersIn, session: Session = Depends(get_session)):
    answers_json = json.dumps(payload.answers, ensure_ascii=False)
    profile = session.exec(select(UserProfile).where(UserProfile.user_id == user_id)).first()
    if profile:
        profile.answers_json = answers_json
        profile.profile_summary = None
        profile.updated_at = datetime.utcnow()
    else:
        profile = UserProfile(user_id=user_id, answers_json=answers_json, updated_at=datetime.utcnow())
        session.add(profile)
    session.commit()
    session.refresh(profile)
    STAGE_CACHE.clear()
    _sync_stage_profile(session, user_id)
    return {"ok": True}


@app.post("/users/{user_id}/personality/generate-summary")
def generate_personality_summary(user_id: int, session: Session = Depends(get_session)):
    profile = session.exec(select(UserProfile).where(UserProfile.user_id == user_id)).first()
    if not profile:
        raise HTTPException(404, "Personality answers not found")
    answers = json.loads(profile.answers_json or "{}")
    if not answers:
        raise HTTPException(400, "No personality answers provided yet")
    content = "\n".join(f"{q}: {a}" for q, a in answers.items())
    profile.profile_summary = llm_text(PERSONALITY_SUMMARY_PROMPT, content, temperature=0.2).strip()
    profile.updated_at = datetime.utcnow()
    session.add(profile)
    session.commit()
    STAGE_CACHE.clear()
    _sync_stage_profile(session, user_id)
    return {"summary": profile.profile_summary}


@app.post("/analyze-text", response_model=TextAnalysisResponse)
def analyze_text(payload: TextAnalysisRequest, session: Session = Depends(get_session)):
    reply_language = payload.reply_language or (payload.meta or {}).get("reply_language")
    pipeline_result = _run_stage_pipeline(session, payload.user_id, payload.text, reply_language)
    item, _, _ = _save_item_from_pipeline(
        session,
        user_id=payload.user_id,
        source="text",
        raw_text=payload.text,
        transcript=None,
        pipeline_result=pipeline_result,
        reply_language=reply_language,
    )
    return item_to_response(item)


@app.post("/transcribe-and-analyze")
async def transcribe_and_analyze(
    user_id: Optional[int] = None,
    reply_language: Optional[str] = None,
    file: UploadFile = File(...),
    session: Session = Depends(get_session),
):
    suffix = os.path.splitext(file.filename)[-1] or ".m4a"
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name

    try:
        transcript_text = _transcribe_audio_file(tmp_path)
        pipeline_result = _run_stage_pipeline(session, user_id, transcript_text, reply_language)
        item, meta, normalized_pipeline = _save_item_from_pipeline(
            session,
            user_id=user_id,
            source="voice",
            raw_text=transcript_text,
            transcript=transcript_text,
            pipeline_result=pipeline_result,
            reply_language=reply_language,
        )
        response = item_to_response(item).model_dump()
        response["assistant"] = {
            "text": item.details or transcript_text,
            "english": normalized_pipeline.get("remodeled_english", ""),
            "tamil": normalized_pipeline.get("tamil_text", ""),
            "theni_tamil": normalized_pipeline.get("theni_tamil_text", ""),
        }
        response["pipeline"] = normalized_pipeline
        response["meta"] = meta
        return response
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass


@app.post("/api/transcribe-and-analyze")
async def api_transcribe_and_analyze(
    user_id: Optional[int] = None,
    reply_language: Optional[str] = None,
    file: UploadFile = File(...),
    session: Session = Depends(get_session),
):
    suffix = os.path.splitext(file.filename)[-1] or ".m4a"
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name

    try:
        transcript_text = _transcribe_audio_file(tmp_path)
        pipeline_result = _run_stage_pipeline(session, user_id, transcript_text, reply_language)
        item, meta, normalized_pipeline = _save_item_from_pipeline(
            session,
            user_id=user_id,
            source="voice",
            raw_text=transcript_text,
            transcript=transcript_text,
            pipeline_result=pipeline_result,
            reply_language=reply_language,
        )
        return _build_chat_response(item, meta, normalized_pipeline)
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass


@app.get("/items", response_model=List[TextAnalysisResponse])
def list_items(session: Session = Depends(get_session), user_id: Optional[int] = None):
    query = select(Item).order_by(Item.created_at.desc())
    if user_id is not None:
        query = query.where(Item.user_id == user_id)
    items = session.exec(query).all()
    return [item_to_response(i) for i in items]


@app.get("/items/{item_id}", response_model=TextAnalysisResponse)
def get_item(item_id: int, session: Session = Depends(get_session)):
    item = session.get(Item, item_id)
    if not item:
        raise HTTPException(404, "Item not found")
    return item_to_response(item)


DOCS_BASE_DIR = "generated_docs"
PDF_BASE_DIR = os.path.join(DOCS_BASE_DIR, "pdf")
EXCEL_BASE_DIR = os.path.join(DOCS_BASE_DIR, "excel")
PPT_BASE_DIR = os.path.join(DOCS_BASE_DIR, "ppt")
DOCX_BASE_DIR = os.path.join(DOCS_BASE_DIR, "docx")


def ensure_dir(path: str):
    os.makedirs(path, exist_ok=True)


def generate_docx(item: Item) -> str:
    from docx import Document

    ensure_dir(DOCX_BASE_DIR)
    cat = os.path.join(DOCX_BASE_DIR, item.category or "Other")
    ensure_dir(cat)
    path = os.path.join(cat, f"item_{item.id}.docx")
    doc = Document()
    doc.add_heading(item.title or f"Item {item.id}", level=1)
    doc.add_paragraph(f"Intent: {item.intent}")
    doc.add_paragraph(f"Category: {item.category}")
    if item.datetime_str:
        doc.add_paragraph(f"When: {item.datetime_str}")
    doc.add_paragraph(item.details or item.raw_text)
    doc.save(path)
    return path


def generate_pdf(item: Item) -> str:
    from fpdf import FPDF

    ensure_dir(PDF_BASE_DIR)
    cat = os.path.join(PDF_BASE_DIR, item.category or "Other")
    ensure_dir(cat)
    path = os.path.join(cat, f"item_{item.id}.pdf")

    pdf = FPDF()
    pdf.add_page()

    font_candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    ]
    font_path = next((p for p in font_candidates if os.path.exists(p)), None)

    if font_path:
        pdf.add_font("DejaVu", "", font_path, uni=True)
        pdf.set_font("DejaVu", size=14)
    else:
        pdf.set_font("Arial", size=12)

    pdf.cell(0, 10, txt=item.title or f"Item {item.id}", ln=True)
    pdf.set_font_size(11)
    pdf.multi_cell(0, 8, txt=f"Intent: {item.intent}")
    pdf.multi_cell(0, 8, txt=f"Category: {item.category}")
    if item.datetime_str:
        pdf.multi_cell(0, 8, txt=f"When: {item.datetime_str}")
    pdf.ln(2)
    pdf.multi_cell(0, 8, txt=item.details or item.raw_text)

    pdf.output(path)
    return path


def generate_excel(item: Item) -> str:
    from openpyxl import Workbook

    ensure_dir(EXCEL_BASE_DIR)
    cat = os.path.join(EXCEL_BASE_DIR, item.category or "Other")
    ensure_dir(cat)
    path = os.path.join(cat, f"item_{item.id}.xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Item"
    rows = [
        ("ID", item.id),
        ("Title", item.title),
        ("Intent", item.intent),
        ("Category", item.category),
        ("When", item.datetime_str),
        ("Details", item.details or item.raw_text),
    ]
    for i, (k, v) in enumerate(rows, start=1):
        ws.cell(row=i, column=1, value=k)
        ws.cell(row=i, column=2, value=v)
    wb.save(path)
    return path


def generate_ppt(item: Item) -> str:
    from pptx import Presentation

    ensure_dir(PPT_BASE_DIR)
    cat = os.path.join(PPT_BASE_DIR, item.category or "Other")
    ensure_dir(cat)
    path = os.path.join(cat, f"item_{item.id}.pptx")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = item.title or f"Item {item.id}"
    tf = slide.placeholders[1].text_frame
    tf.text = f"Intent: {item.intent}\nCategory: {item.category}"
    if item.datetime_str:
        tf.add_paragraph().text = f"When: {item.datetime_str}"
    tf.add_paragraph().text = item.details or item.raw_text

    prs.save(path)
    return path


@app.post("/items/{item_id}/generate-pdf")
def item_generate_pdf(item_id: int, session: Session = Depends(get_session)):
    item = session.get(Item, item_id)
    if not item:
        raise HTTPException(404, "Item not found")
    path = generate_pdf(item)
    return {"ok": True, "path": path, "download_url": f"/download?path={path}"}


@app.post("/items/{item_id}/generate-excel")
def item_generate_excel(item_id: int, session: Session = Depends(get_session)):
    item = session.get(Item, item_id)
    if not item:
        raise HTTPException(404, "Item not found")
    path = generate_excel(item)
    return {"ok": True, "path": path, "download_url": f"/download?path={path}"}


@app.post("/items/{item_id}/generate-ppt")
def item_generate_ppt(item_id: int, session: Session = Depends(get_session)):
    item = session.get(Item, item_id)
    if not item:
        raise HTTPException(404, "Item not found")
    path = generate_ppt(item)
    return {"ok": True, "path": path, "download_url": f"/download?path={path}"}


@app.post("/items/{item_id}/generate-docx")
def item_generate_docx(item_id: int, session: Session = Depends(get_session)):
    item = session.get(Item, item_id)
    if not item:
        raise HTTPException(404, "Item not found")
    path = generate_docx(item)
    return {"ok": True, "path": path, "download_url": f"/download?path={path}"}


@app.get("/download")
def download_generated(path: str):
    if not os.path.isfile(path):
        raise HTTPException(404, "File not found")
    filename = os.path.basename(path)
    return FileResponse(path, filename=filename)